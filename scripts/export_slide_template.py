#!/usr/bin/env python3
"""
Extract slide layout information from a PowerPoint deck and emit a template JSON skeleton.

Example:
    python scripts/export_slide_template.py /path/to/Cover_Page.pptx --slide 0 --template-id cover_page
"""

import argparse
import json
import os
from typing import Any, Dict, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def round_inches(value) -> float:
    """Round a python-pptx length to inches with three decimal places."""
    return round(value.inches, 3)


def sanitize_role(name: str, fallback: str) -> str:
    """Return a lowercase snake-case role name."""
    base = (name or fallback).strip().lower()
    for ch in (" ", "-", ".", "(", ")", "[", "]", "{", "}", "#", "%"):
        base = base.replace(ch, "_")
    base = "".join(ch for ch in base if ch.isalnum() or ch == "_")
    while "__" in base:
        base = base.replace("__", "_")
    base = base.strip("_")
    return base or fallback


def shape_role(shape, idx: int) -> str:
    """Best-effort role detection for a shape."""
    name = getattr(shape, "name", "") or ""
    lowered = name.lower()
    if "title" in lowered:
        return "title"
    if "subtitle" in lowered:
        return "subtitle"
    if "content" in lowered:
        return "content"
    if "picture" in lowered or "image" in lowered:
        return "hero_image"

    if getattr(shape, "placeholder_format", None):
        placeholder_type = getattr(shape.placeholder_format, "type", None)
        if placeholder_type:
            return sanitize_role(str(placeholder_type).split(".")[-1], f"shape_{idx}")
    return sanitize_role(name, f"shape_{idx}")


def shape_to_element(shape, idx: int) -> Dict[str, Any]:
    """Convert a pptx shape into template element metadata."""
    element: Dict[str, Any] = {
        "type": "shape",
        "role": shape_role(shape, idx),
        "position": {
            "left": round_inches(shape.left),
            "top": round_inches(shape.top),
            "width": round_inches(shape.width),
            "height": round_inches(shape.height),
        },
    }

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        element["type"] = "image"
        element["effects"] = ["professional_shadow"]
        return element

    if shape.has_table:
        element["type"] = "table"
        return element

    if shape.has_text_frame:
        element["type"] = "text"
        text = shape.text.strip()
        if text:
            element["placeholder_text"] = text

        styling: Dict[str, Any] = {}
        # Heuristics for common placeholder roles
        role = element["role"]
        if "title" in role:
            styling.update({"font_type": "title", "font_size": "large"})
        elif "subtitle" in role:
            styling.update({"font_type": "subtitle", "font_size": "medium"})
        else:
            styling.update({"font_type": "body", "font_size": "medium"})

        element["styling"] = styling
        return element

    return element


def export_template(args: argparse.Namespace) -> Dict[str, Any]:
    """Build template dictionary for the requested slide."""
    if not os.path.exists(args.presentation):
        raise FileNotFoundError(f"Presentation not found: {args.presentation}")

    prs = Presentation(args.presentation)
    slide_count = len(prs.slides)
    if slide_count == 0:
        raise ValueError("Presentation contains no slides")
    if args.slide < 0 or args.slide >= slide_count:
        raise IndexError(f"Slide index {args.slide} out of range (0-{slide_count - 1})")

    slide = prs.slides[args.slide]
    elements: List[Dict[str, Any]] = []
    seen_roles: Dict[str, int] = {}
    for idx, shape in enumerate(slide.shapes, start=1):
        element = shape_to_element(shape, idx)
        role = element["role"]
        counter = seen_roles.get(role, 0)
        if counter:
            element["role"] = f"{role}_{counter+1}"
        seen_roles[role] = counter + 1
        elements.append(element)

    template: Dict[str, Any] = {
        "template_id": args.template_id or sanitize_role(os.path.splitext(os.path.basename(args.presentation))[0], "custom_slide"),
        "name": args.name or os.path.splitext(os.path.basename(args.presentation))[0],
        "description": args.description or f"Auto-extracted from {os.path.basename(args.presentation)}",
        "layout_type": args.layout_type,
        "typography_style": args.typography_style,
        "elements": elements,
        "background": {
            "type": "solid",
            "color_role": args.background_color_role,
        },
    }

    return template


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Export slide geometry as MCP template JSON.")
    parser.add_argument("presentation", help="Path to .pptx file containing the slide to export")
    parser.add_argument("--slide", type=int, default=0, help="Zero-based slide index to export (default: 0)")
    parser.add_argument("--template-id", help="Template ID to use in the JSON output")
    parser.add_argument("--name", help="Human-friendly template name")
    parser.add_argument("--description", help="Template description text")
    parser.add_argument("--layout-type", default="title", help="Layout type (e.g. title, content, custom)")
    parser.add_argument("--typography-style", default="modern_sans", help="Typography style key (default: modern_sans)")
    parser.add_argument(
        "--background-color-role",
        default="light",
        help="Color role key for the background (default: light)",
    )
    parser.add_argument("--output", "-o", help="Write JSON to file instead of stdout")
    return parser


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()
    template = export_template(args)
    output = json.dumps(template, indent=2)

    if args.output:
        with open(args.output, "w", encoding="utf-8") as fh:
            fh.write(output)
    else:
        print(output)


if __name__ == "__main__":
    main()

